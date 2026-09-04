import os
import asyncio
import re
import copy
import subprocess
import tempfile
import uuid
import html
import certifi
from aiohttp import web
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import CommandStart, Command, StateFilter
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.client.default import DefaultBotProperties
from aiogram.enums import ParseMode, ChatAction
from aiogram.types import InlineKeyboardMarkup, InlineKeyboardButton, CallbackQuery, BotCommand

from pptx import Presentation
from docxtpl import DocxTemplate, RichText
from pymongo import MongoClient

# ==================== CONFIGURATION ====================
API_TOKEN = os.getenv("BOT_TOKEN", "8941822350:AAF2H5oiIdvQ70t1TFhX8lhkvHlZLzHBkPc")

MONGO_URI = os.getenv(
    "MONGO_URI", 
    "mongodb+srv://mpcpawan:RswOqZ4uy3UQtM3Q@cluster0.edkvmpu.mongodb.net/"
)

BASE_DIR = os.path.dirname(__file__)
PPT_TEMPLATE = os.path.join(BASE_DIR, "template.pptx")
DOCX_TEMPLATE = os.path.join(BASE_DIR, "template.docx")

client = MongoClient(MONGO_URI, tlsCAFile=certifi.where(), serverSelectionTimeoutMS=5000)
db = client["rcc_quiz_db"]
tests_col = db["test_papers"]

bot = Bot(token=API_TOKEN, default=DefaultBotProperties(parse_mode=ParseMode.HTML))
dp = Dispatcher(storage=MemoryStorage())

class QuizForm(StatesGroup):
    waiting_for_topic = State()
    waiting_for_format = State()
    collecting_questions = State()


# ==================== SET TELEGRAM MENU COMMANDS ====================

async def setup_bot_commands(bot_instance: Bot):
    commands = [
        BotCommand(command="start", description="🤖 बॉट शुरू करें"),
        BotCommand(command="create", description="📝 नया टेस्ट बनाएं"),
        BotCommand(command="prompt", description="✨ Gemini AI Prompt (फोटो/PDF से प्रश्न बनाएं)"),
        BotCommand(command="cancel", description="❌ चालू प्रक्रिया रद्द करें"),
        BotCommand(command="help", description="❓ सहायता एवं निर्देश"),
        BotCommand(command="mytests", description="📂 मेरे बनाए गए हालिया टेस्ट"),
        BotCommand(command="ppt", description="📊 PPT PDF बनाएं (/ppt ID)"),
        BotCommand(command="test", description="📄 Test PDF बनाएं (/test ID)"),
        BotCommand(command="answer", description="✅ Answer PDF बनाएं (/answer ID)"),
        BotCommand(command="stats", description="📈 कुल टेस्ट के आंकड़े"),
    ]
    await bot_instance.set_my_commands(commands)


# ==================== HELPER FUNCTIONS ====================

def convert_to_pdf(input_file, output_dir="."):
    abs_input = os.path.abspath(input_file)
    out_name = os.path.basename(input_file).rsplit('.', 1)[0] + '.pdf'
    abs_output = os.path.abspath(os.path.join(output_dir, out_name))

    if os.name == 'nt':
        libre_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "soffice"
        ]
        for soffice_path in libre_paths:
            try:
                cmd = [soffice_path, "--headless", "--convert-to", "pdf", abs_input, "--outdir", output_dir]
                subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                if os.path.exists(abs_output):
                    return
            except Exception:
                continue

        if input_file.endswith('.pptx') or input_file.endswith('.ppt'):
            try:
                import win32com.client
                try:
                    import pythoncom
                    pythoncom.CoInitialize()
                except Exception:
                    pass
                
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                deck = powerpoint.Presentations.Open(abs_input, WithWindow=False)
                deck.SaveAs(abs_output, 32)
                deck.Close()
                return
            except Exception:
                raise Exception("Windows पर PPT को PDF बनाने के लिए MS PowerPoint का होना आवश्यक है!")

        if input_file.endswith('.docx'):
            try:
                from docx2pdf import convert
                convert(abs_input, abs_output)
                return
            except Exception as e:
                raise Exception(f"DOCX to PDF Error: {str(e)}")

        raise Exception("PDF कनवर्टर फ़ाइल जनरेट करने में असमर्थ रहा।")
    else:
        cmd = ["libreoffice", "--headless", "--convert-to", "pdf", abs_input, "--outdir", output_dir]
        subprocess.run(cmd, check=True)

def parse_raw_text(raw_text):
    questions_list = []
    q_blocks = re.split(r'\n(?=\s*\d+[\.\)\-])', '\n' + raw_text.strip())
    
    for block in q_blocks:
        if not block.strip(): 
            continue
        lines = [line.strip() for line in block.split('\n') if line.strip()]
        if not lines: 
            continue
            
        q_text = re.sub(r'^\d+[\.\)\-]\s*', '', lines[0])
        full_text = "\n".join(lines[1:])
        
        opt_a = re.search(r'(?:^|\n|\s*)(?:\(a\)|a[\.\)\-])\s*(.*?)(?=(?:\(b\)|b[\.\)\-])|$)', full_text, re.DOTALL | re.IGNORECASE)
        opt_b = re.search(r'(?:^|\n|\s*)(?:\(b\)|b[\.\)\-])\s*(.*?)(?=(?:\(c\)|c[\.\)\-])|$)', full_text, re.DOTALL | re.IGNORECASE)
        opt_c = re.search(r'(?:^|\n|\s*)(?:\(c\)|c[\.\)\-])\s*(.*?)(?=(?:\(d\)|d[\.\)\-])|$)', full_text, re.DOTALL | re.IGNORECASE)
        opt_d = re.search(r'(?:^|\n|\s*)(?:\(d\)|d[\.\)\-])\s*(.*?)(?=$)', full_text, re.DOTALL | re.IGNORECASE)
        
        questions_list.append({
            'text': q_text.strip(),
            'a': opt_a.group(1).strip() if opt_a else '',
            'b': opt_b.group(1).strip() if opt_b else '',
            'c': opt_c.group(1).strip() if opt_c else '',
            'd': opt_d.group(1).strip() if opt_d else ''
        })
    return questions_list

def format_docx_option(label, opt_text, show_answer=False):
    if not opt_text: 
        return ""
    rt = RichText()
    is_answer = "✅" in opt_text or "*" in opt_text
    cleaned = opt_text.replace("✅", "").replace("*", "").strip()
    
    if show_answer and is_answer:
        rt.add(f"{label} {cleaned}", bold=True)
    else:
        rt.add(f"{label} {cleaned}", bold=False)
    return rt


# ==================== PDF GENERATION ENGINE ====================

async def generate_and_send(chat_id, doc_id, gen_type):
    try:
        row = tests_col.find_one({"_id": doc_id})
    except Exception as e:
        await bot.send_message(chat_id, f"❌ <b>Database Error:</b> {html.escape(str(e))}")
        return

    if not row:
        await bot.send_message(chat_id, "❌ <b>ID नहीं मिला!</b> कृपया सही ID दर्ज करें।")
        return

    topic = row["topic"]
    raw_text = row["raw_text"]
    parsed_qs = parse_raw_text(raw_text)
    
    await bot.send_chat_action(chat_id=chat_id, action=ChatAction.TYPING)
    msg = await bot.send_message(chat_id, f"⏳ <b>{gen_type}</b> जनरेट हो रहा है, कृपया प्रतीक्षा करें...")

    temp_dir = tempfile.gettempdir()
    output_pdf = os.path.join(temp_dir, f"{topic.replace(' ', '_')}_{gen_type}_{doc_id}.pdf")
    
    try:
        loop = asyncio.get_running_loop()

        if gen_type == "PPT":
            if not os.path.exists(PPT_TEMPLATE):
                await msg.edit_text("❌ <b>Template Missing:</b> `template.pptx` नहीं मिला!")
                return
                
            output_file = os.path.join(temp_dir, f"temp_{doc_id}.pptx")
            prs = Presentation(PPT_TEMPLATE)
            base_slide = prs.slides[0]
            base_shapes_elements = [copy.deepcopy(shape.element) for shape in base_slide.shapes]
            
            for index, q in enumerate(parsed_qs, 1):
                cl_a = q['a'].replace("✅", "").replace("*", "").strip()
                cl_b = q['b'].replace("✅", "").replace("*", "").strip()
                cl_c = q['c'].replace("✅", "").replace("*", "").strip()
                cl_d = q['d'].replace("✅", "").replace("*", "").strip()
                
                q_text_val = f"Q{index}. {q['text']}"
                a_val = f"A) {cl_a}"
                b_val = f"B) {cl_b}"
                c_val = f"C) {cl_c}"
                d_val = f"D) {cl_d}"
                
                replacements = {
                    '{{TOPIC}}': topic, 
                    '{{QUESTION}}': q_text_val,
                    '{{OPTION_A}}': a_val, 
                    '{{OPTION_B}}': b_val,
                    '{{OPTION_C}}': c_val, 
                    '{{OPTION_D}}': d_val
                }
                
                target_slide = base_slide if index == 1 else prs.slides.add_slide(prs.slide_layouts[6])
                if index != 1:
                    for el in base_shapes_elements:
                        target_slide.shapes._spTree.insert_element_before(copy.deepcopy(el), 'p:extLst')
                
                for shape in target_slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            full_p_text = "".join(run.text for run in paragraph.runs)
                            changed = False
                            for key, val in replacements.items():
                                if key in full_p_text:
                                    full_p_text = full_p_text.replace(key, str(val))
                                    changed = True
                            if changed and paragraph.runs:
                                paragraph.runs[0].text = full_p_text
                                for run in paragraph.runs[1:]:
                                    run.text = ""

            prs.save(output_file)
            await loop.run_in_executor(None, convert_to_pdf, output_file, temp_dir)

        else: # DOCX Formats (Test PDF & Answer Test PDF)
            if not os.path.exists(DOCX_TEMPLATE):
                await msg.edit_text("❌ <b>Template Missing:</b> `template.docx` नहीं मिला!")
                return
                
            output_file = os.path.join(temp_dir, f"temp_{doc_id}.docx")
            
            # फिक्स: फाइल पाथ का उपयोग करके 'seek of closed file' एरर पूरी तरह खत्म
            doc = DocxTemplate(DOCX_TEMPLATE)
            
            show_answers = (gen_type == "Answer Test PDF")
            
            formatted_qs = []
            for i, q in enumerate(parsed_qs, 1):
                formatted_qs.append({
                    'id': i, 
                    'text': q['text'],
                    'a': format_docx_option("(a)", q['a'], show_answers),
                    'b': format_docx_option("(b)", q['b'], show_answers),
                    'c': format_docx_option("(c)", q['c'], show_answers),
                    'd': format_docx_option("(d)", q['d'], show_answers),
                })
            
            context = {'topic_name': topic, 'questions': formatted_qs}
            doc.render(context)
            doc.save(output_file)
            await loop.run_in_executor(None, convert_to_pdf, output_file, temp_dir)

        generated_pdf_path = output_file.rsplit('.', 1)[0] + ".pdf"
        if os.path.exists(generated_pdf_path):
            os.rename(generated_pdf_path, output_pdf)
            await bot.send_chat_action(chat_id=chat_id, action=ChatAction.UPLOAD_DOCUMENT)
            await bot.send_document(
                chat_id, types.FSInputFile(output_pdf),
                caption=f"📄 आपका <b>{gen_type}</b> तैयार है!\n🆔 <b>ID:</b> <code>{doc_id}</code>"
            )
        else:
            await msg.edit_text("❌ PDF जनरेट करने में विफलता हुई।")
        
        if os.path.exists(output_file): 
            os.remove(output_file)
        if os.path.exists(output_pdf): 
            os.remove(output_pdf)
        await msg.delete()

    except Exception as e:
        await msg.edit_text(f"❌ <b>Error:</b> {html.escape(str(e))}")


# ==================== GLOBAL COMMAND HANDLERS ====================

@dp.message(CommandStart(), StateFilter("*"))
async def cmd_start(message: types.Message, state: FSMContext):
    await state.clear()
    welcome_text = (
        "🤖 <b>Rajesh Competition Centre Bot में आपका स्वागत है!</b>\n\n"
        "यह बॉट प्रतियोगी परीक्षाओं के लिए क्विज, टेस्ट पेपर्स और PPT PDFs जनरेट करता है।\n\n"
        "📌 <b>उपलब्ध मुख्य कमांड्स:</b>\n"
        "• /create - नया टेस्ट बनाएं\n"
        "• /prompt - फोटो / PDF से AI प्रश्न बनाने का Prompt\n"
        "• /mytests - अपने हालिया टेस्ट देखें\n"
        "• /cancel - प्रक्रिया रद्द करें\n"
        "• /help - प्रयोग करने की गाइड\n"
        "• /stats - डेटाबेस के आंकड़े\n\n"
        "शुरू करने के लिए नीचे <b>'Create'</b> बटन दबाएं 👇"
    )
    keyboard = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📝 Create (नया बनाएं)", callback_data="btn_create")],
        [InlineKeyboardButton(text="✨ AI Prompt (फोटो / PDF से प्रश्न)", callback_data="btn_prompt")],
        [InlineKeyboardButton(text="❓ सहायता (Help)", callback_data="btn_help")]
    ])
    await message.reply(welcome_text, reply_markup=keyboard)

@dp.message(Command("cancel"), StateFilter("*"))
async def cmd_cancel(message: types.Message, state: FSMContext):
    current_state = await state.get_state()
    if current_state is None:
        await message.reply("ℹ️ कोई सक्रिय प्रक्रिया जारी नहीं है।")
        return
    
    await state.clear()
    await message.reply("❌ <b>चालू प्रक्रिया सफलतापूर्वक रद्द (Cancel) कर दी गई है!</b>\n\nनया टेस्ट बनाने के लिए /create टाइप करें।")

@dp.message(Command("create"), StateFilter("*"))
async def cmd_create(message: types.Message, state: FSMContext):
    await state.clear()
    await state.set_state(QuizForm.waiting_for_topic)
    await message.reply("🎯 <b>कृपया अपने टेस्ट पेपर का नाम (Topic Name) दर्ज करें:</b>\n\n<i>(उदाहरण: हर्यक वंश एवं मगध साम्राज्य)</i>")

@dp.message(Command("prompt"), StateFilter("*"))
@dp.callback_query(F.data == "btn_prompt")
async def cmd_prompt(event: types.Message | CallbackQuery):
    prompt_text = (
        "✨ <b>Gemini AI क्विज़ मेकर प्रॉम्प्ट (Master Prompt)</b>\n\n"
        "अगर आपके पास किसी <b>किताब की फोटो, AI, PYQ, PDF या थ्योरी नोट्स</b> हैं, या फिर आप टॉपिक के नाम से प्रश्न बनवाना चाहते हैं, तो नीचे दिए गए बॉक्स पर क्लिक करके <b>Copy</b> करें:\n\n"
        "<code>विषय/टॉपिक (Topic Name): [यहाँ अपना टॉपिक लिखें या खाली छोड़ें अगर PDF अटैच है]\n"
        "प्रश्नों की संख्या (Total Questions): [जितने प्रश्न चाहिए जैसे 30, 50 लिखें]\n\n"
        "कृपया दिए गए Photo / PDF / PYQ / थ्योरी नोट्स को ध्यान से पढ़ें और इनसे ऑब्जेक्टिव प्रश्न (MCQs) बनाकर मुझे बिल्कुल इसी फॉर्मेट में दें:\n\n"
        "1. भारत की राजधानी क्या है?\n"
        "a) मुंबई\n"
        "b) नई दिल्ली ✅\n"
        "c) कोलकाता\n"
        "d) चेन्नई\n\n"
        "नियम (Strict Rules):\n"
        "1. जो उत्तर सही (Correct Answer) हो, उसके विकल्प के अंत में अनिवार्य रूप से '✅' ग्रीन टिक लगाएं।\n"
        "2. अगर यह थ्योरी नोट्स हैं, तो उससे सबसे महत्वपूर्ण प्रश्न खुद बनाएं।\n"
        "3. अगर यह फोटो या PDF है, तो उसमें मौजूद सभी प्रश्नों को ऊपर दिए गए फॉर्मेट में डिजिटल टेक्स्ट में बदलें।\n"
        "4. उत्तर देने में कोई भी फालतू बात या परिचय न लिखें। सिर्फ और सिर्फ प्रश्नों की लिस्ट दें।\n"
        "5. आपके द्वारा दिए जाने वाले सारे के सारे प्रश्न एक ही सिंगल कोडिंग बॉक्स (Code Block) के अंदर होने चाहिए, ताकि एक क्लिक में पूरा टेक्स्ट कॉपी किया जा सके।</code>\n\n"
        "📌 <b>उपयोग करने की विधि:</b>\n"
        "1. ऊपर दिए गए कोड पर टैप करके <b>Copy</b> करें।\n"
        "2. <b>Google Gemini App</b> (या चैट) में जाएं।\n"
        "3. अपनी फोटो / PDF अटैच करें (या बिना अटैच किए टॉपिक और प्रश्नों की संख्या भरें) और यह Prompt पेस्ट करके भेजें।\n"
        "4. Gemini से मिले उत्तर को सीधे कॉपी करके यहाँ बॉट में <b>/create</b> दबाकर पेस्ट कर दें!"
    )
    if isinstance(event, CallbackQuery):
        await event.message.reply(prompt_text)
        await event.answer()
    else:
        await event.reply(prompt_text)

@dp.message(Command("help"), StateFilter("*"))
@dp.callback_query(F.data == "btn_help")
async def cmd_help(event: types.Message | CallbackQuery):
    help_text = (
        "📖 <b>RCC Quiz Bot - गाइड एवं सहायता</b>\n\n"
        "<b>1. प्रश्न कैसे भेजें?</b>\n"
        "प्रश्न निम्न फॉर्मेट में भेजें:\n"
        "<code>1. भारत की राजधानी क्या है?\n"
        "a) मुंबई\n"
        "b) नई दिल्ली ✅\n"
        "c) कोलकाता\n"
        "d) चेन्नई</code>\n\n"
        "<b>2. फोटो / PDF से प्रश्न कैसे बनाएं?</b>\n"
        "• /prompt कमांड टाइप करें और दिए गए प्रॉम्प्ट को Gemini AI में उपयोग करें।\n\n"
        "<b>3. ID से पुनः PDF डाउनलोड करना:</b>\n"
        "• PPT के लिए: <code>/ppt &lt;ID&gt;</code>\n"
        "• केवल प्रश्नों के लिए: <code>/test &lt;ID&gt;</code>\n"
        "• उत्तर कुंजी सहित: <code>/answer &lt;ID&gt;</code>\n\n"
        "<b>4. प्रक्रिया रद्द करना:</b>\n"
        "किसी भी समय /cancel भेजकर प्रक्रिया रोक सकते हैं।"
    )
    if isinstance(event, CallbackQuery):
        await event.message.reply(help_text)
        await event.answer()
    else:
        await event.reply(help_text)

@dp.message(Command("mytests"), StateFilter("*"))
async def cmd_mytests(message: types.Message):
    try:
        records = list(tests_col.find().sort("_id", -1).limit(5))
        if not records:
            await message.reply("📂 अभी तक कोई टेस्ट पेपर सेव नहीं हुआ है।")
            return
            
        text = "📂 <b>हाल ही में बनाए गए टेस्ट पेपर्स:</b>\n\n"
        for r in records:
            text += f"🆔 <b>ID:</b> <code>{r['_id']}</code> | 📌 <b>Topic:</b> {r.get('topic', 'N/A')}\n"
            text += f"└ <code>/test {r['_id']}</code> | <code>/answer {r['_id']}</code> | <code>/ppt {r['_id']}</code>\n\n"
            
        await message.reply(text)
    except Exception as e:
        await message.reply(f"❌ <b>Error:</b> {html.escape(str(e))}")

@dp.message(Command("stats"), StateFilter("*"))
async def cmd_stats(message: types.Message):
    try:
        count = tests_col.count_documents({})
        await message.reply(f"📊 <b>डेटाबेस आंकड़े:</b>\n\nकुल सेव किए गए टेस्ट पेपर्स: <b>{count}</b>")
    except Exception as e:
        await message.reply(f"❌ <b>Error:</b> {html.escape(str(e))}")


# ==================== CREATION FLOW HANDLERS ====================

@dp.callback_query(F.data == "btn_create")
async def ask_topic(callback: CallbackQuery, state: FSMContext):
    await state.clear()
    await state.set_state(QuizForm.waiting_for_topic)
    await callback.message.edit_text("🎯 <b>कृपया अपने टेस्ट पेपर का नाम (Topic Name) दर्ज करें:</b>")

@dp.message(QuizForm.waiting_for_topic)
async def process_topic(message: types.Message, state: FSMContext):
    await state.update_data(topic=message.text.strip())
    
    keyboard = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📊 PPT (पीपीटी)", callback_data="fmt_PPT")],
        [InlineKeyboardButton(text="📄 Test PDF (टेस्ट पीडीएफ)", callback_data="fmt_Test PDF")],
        [InlineKeyboardButton(text="✅ Answer Test PDF (आंसर सहित)", callback_data="fmt_Answer Test PDF")]
    ])
    await message.reply("📝 <b>टॉपिक सेव हो गया!</b>\n\nअब चुनें कि किस फॉर्मेट में जनरेट करना चाहते हैं:", reply_markup=keyboard)
    await state.set_state(QuizForm.waiting_for_format)

@dp.callback_query(QuizForm.waiting_for_format)
async def ask_questions(callback: CallbackQuery, state: FSMContext):
    fmt = callback.data.replace("fmt_", "")
    await state.update_data(selected_format=fmt, raw_questions="")
    
    sample = (
        "1. भारत की राजधानी क्या है?\n"
        "a) मुंबई\n"
        "b) नई दिल्ली ✅\n"
        "c) कोलकाता\n"
        "d) चेन्नई"
    )
    
    msg = (
        f"✅ आपने <b>{fmt}</b> चुना है।\n\n"
        "👇 <b>कृपया अपने प्रश्न इस फॉर्मेट में भेजें:</b>\n"
        f"<code>{sample}</code>\n\n"
        "📌 <i>नोट: प्रश्न भेजने के बाद चाहें तो और भेजें, समाप्त होने पर <b>/done</b> टाइप करें।\n"
        "रद्द करने के लिए <b>/cancel</b> दबाएं।</i>"
    )
    await callback.message.edit_text(msg)
    await state.set_state(QuizForm.collecting_questions)

@dp.message(QuizForm.collecting_questions)
async def collect_questions(message: types.Message, state: FSMContext):
    text = message.text.strip()
    
    if text.startswith('/') and text.lower() != '/done':
        return

    if text.lower() == '/done':
        user_data = await state.get_data()
        raw_text = user_data.get('raw_questions', '')
        topic = user_data.get('topic', 'Test')
        selected_format = user_data.get('selected_format', 'Test PDF')
        
        if not raw_text.strip():
            await message.reply("❌ कोई प्रश्न नहीं मिला! कृपया पहले प्रश्न भेजें।")
            return

        doc_id = uuid.uuid4().hex[:6].upper()
        
        try:
            tests_col.insert_one({
                "_id": doc_id,
                "topic": topic,
                "raw_text": raw_text
            })
        except Exception as e:
            await message.reply(f"❌ <b>Database Insert Failed:</b> {html.escape(str(e))}")
            return

        success_msg = (
            f"✅ <b>डेटा सेव हो गया!</b>\n\n"
            f"🔑 <b>आपकी Test ID:</b> <code>{doc_id}</code>\n\n"
            f"💡 <b>भविष्य में इस ID से डाउनलोड करें:</b>\n"
            f"• PPT: <code>/ppt {doc_id}</code>\n"
            f"• Test PDF: <code>/test {doc_id}</code>\n"
            f"• Answer PDF: <code>/answer {doc_id}</code>\n\n"
            f"<i>अभी आपका ({selected_format}) जनरेट किया जा रहा है...</i>"
        )
        await message.reply(success_msg)
        await state.clear()
        
        await generate_and_send(message.chat.id, doc_id, selected_format)
        return

    user_data = await state.get_data()
    current_raw = user_data.get('raw_questions', '')
    new_raw = current_raw + "\n\n" + text if current_raw else text
    await state.update_data(raw_questions=new_raw)
    
    parsed = parse_raw_text(new_raw)
    await message.reply(f"📥 <b>प्रश्न जोड़ दिए गए! (कुल: {len(parsed)})</b>\nऔर प्रश्न भेजें या जनरेट करने के लिए <b>/done</b> टाइप करें।")


# ==================== ID COMMAND HANDLERS ====================

@dp.message(Command("ppt"), StateFilter("*"))
async def cmd_ppt(message: types.Message):
    args = message.text.split()
    if len(args) < 2:
        await message.reply("⚠️ कृपया ID दर्ज करें। उदाहरण: <code>/ppt A1B2C3</code>")
        return
    await generate_and_send(message.chat.id, args[1].upper(), "PPT")

@dp.message(Command("test"), StateFilter("*"))
async def cmd_test(message: types.Message):
    args = message.text.split()
    if len(args) < 2:
        await message.reply("⚠️ कृपया ID दर्ज करें। उदाहरण: <code>/test A1B2C3</code>")
        return
    await generate_and_send(message.chat.id, args[1].upper(), "Test PDF")

@dp.message(Command("answer"), StateFilter("*"))
async def cmd_answer(message: types.Message):
    args = message.text.split()
    if len(args) < 2:
        await message.reply("⚠️ कृपया ID दर्ज करें। उदाहरण: <code>/answer A1B2C3</code>")
        return
    await generate_and_send(message.chat.id, args[1].upper(), "Answer Test PDF")


# ==================== RENDER WEB SERVER ====================

async def handle_ping(request):
    return web.Response(text="RCC Professional Bot is Live!")

async def start_web_server():
    app = web.Application()
    app.router.add_get('/', handle_ping)
    runner = web.AppRunner(app)
    await runner.setup()
    port = int(os.environ.get("PORT", 8080))
    site = web.TCPSite(runner, '0.0.0.0', port)
    await site.start()

async def main():
    # पुरानी वेब सर्विस और पोलिंग संघर्ष रोकने के लिए
    await start_web_server()
    await setup_bot_commands(bot)
    
    # ड्रॉप पेंडिंग अपडेट्स ताकि Conflict Error न आए
    await bot.delete_webhook(drop_pending_updates=True)
    
    print("\n" + "="*50)
    print("🚀 RCC PROFESSIONAL BOT IS LIVE AND RUNNING!")
    print("="*50 + "\n")
    
    await dp.start_polling(bot)

if __name__ == '__main__':
    try:
        asyncio.run(main())
    except (KeyboardInterrupt, SystemExit):
        print("Bot Stopped!")
